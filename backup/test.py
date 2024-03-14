
# LIBRARIES PROJECT

import pyodbc
import pandas as pd
import matplotlib.pyplot as plt
import numpy as np
import seaborn as sns
import textwrap
from datetime import datetime,date
import warnings
from pptx import Presentation 
from pptx.util import Inches
warnings.filterwarnings("ignore")
import os
import glob

 

# GLOBAL VARIABLES

path_dgsante_project = 'C:/Users/amedina/OneDrive - Tecnológica de Alimentos S.A/Proyectos/dgsante'
path_resources = 'resources'
template_pptx = 'template_dgsante.pptx'

inicio_tem_2022_1 = '2022-04-01'
fin_tem_2022_1 = '2022-07-31'

inicio_tem_2022_2 = '2022-11-23'
fin_tem_2022_2 = '2023-01-31'

# PLOT FUNCTIONS

def barplot_df(x,y,df,title):
    ax = sns.barplot(x=x,y=y,data=df)
    ax.bar_label(ax.containers[0])
    ax.get_figure().savefig(path_resources+'/'+title+'.png')
    

def barplot_df_hue(x,y,hue,df,title):
    plt.figure(figsize=(10,10))
    ax = sns.barplot(x=x,y=y,hue=hue,data=df)
    plt.title(title,fontsize=10)
    for i in ax.containers:
        ax.bar_label(i,)
    ax.get_figure().savefig(path_resources+'/'+title+'.png')
    
    
def boxplot_df(df,x,title):
    plt.figure(figsize=(10,10))
    ax = sns.boxplot(x=df[x]).set(title=title)
    ax.get_figure().savefig(path_resources+'/'+title+'.png')
    
    
def piechart_percentage(df,title):
    plt.figure(figsize=(10,10))
    labels = ['Inconforme','Conforme']
    colors = ('#bb2124','#5bc0de')
    plt.title(title,fontsize=10)
    plt.pie(df['Percentage'], labels=labels, colors=colors,autopct='%1.1f%%',
            shadow=True, startangle=90)
    plt.axis('equal')  
    plt.savefig(path_resources+'/'+title+'.png')
    
    
def histplot_df_hue(df,x,hue,title):
    plt.title(title,fontsize=10)
    plt.figure(figsize=(10,10))
    sns.histplot(df, x=x, hue=hue)
    plt.savefig(path_resources+'/'+title+'.png')
    
    
def histplot_df(df,x,title):
    
    plt.title(title,fontsize=10)
    plt.figure(figsize=(10,10))
    sns.histplot(df, x=x,binwidth=3)
    plt.savefig(path_resources+'/'+title+'.png')
   
    
def countplot_percentage_hue(x,hue,df,title):
    plt.figure(figsize=(10,10))
    df = df.groupby(x)[hue].value_counts(normalize=True)
    df = df.mul(100)
    df = df.rename('percent').reset_index()
    
    g = sns.catplot(x=x,y='percent',hue=hue, kind='bar',data=df)
 
    plt.title(title, fontsize=10)
    for p in g.ax.patches:
        percentage = str((p.get_height().round()).astype(int)) + '%'
        x = p.get_x()+ 0.015
        y = p.get_height() * 1.02
        g.ax.annotate(percentage, (x, y),color='black', rotation='horizontal')
    fig = g.ax.get_figure()
    fig.savefig(path_resources+'/'+title+'.png')
  
# DATA PROCESSING FUNCTIONS 

def solve_datetime_sheets(datetime_s):
 
    year = (datetime_s).year 
    mon = (datetime_s).month 
    day = (datetime_s).day 
    hour = (datetime_s).hour
    minute = (datetime_s).minute
    second = (datetime_s).second

    if (( day == 12) & (mon < day)  ):
        return datetime(year, day, mon,hour,minute,second)
    else:
        return datetime(year, mon, day,hour,minute,second)
    
def values_total(x,df):
    all = df[x].value_counts()
    percentage = all/len(df)*100
    df_mv = pd.DataFrame(data=[all,percentage])
    df_mv = df_mv.T
    df_mv.columns = ['Values','Percentage']
    return df_mv.sort_index()    

def identify_outliers(df):
  for i in df.describe().columns:
    Q1 = df.describe().at['25%', i]
    Q3 = df.describe().at['75%', i]
    IQR = Q3-Q1
    LTV = Q1 - 1.5 * IQR 
    UTV = Q3 + 1.5 * IQR 
    print('Nombre de Columna                                         :', i)
    print('Rango Intercuartil IQR                                    :', IQR)
    print('Limite Inferior                                           :', LTV)
    print('Limite Superior                                           :', UTV)
    print('Número de outliers por debajo del bigote inferior         :', df[df[i] < LTV][i].count())
    print('Número de outliers por encima del bigote superior         :', df[df[i] > UTV][i].count())
    print('Total de Ourliers                                         :', (df[df[i] < LTV][i].count()) + (df[df[i] > UTV][i].count()))
    print('\n')
    
def select_values_outliers(df):
  for i in df.describe().columns:
    Q1 = df.describe().at['25%', i]
    Q3 = df.describe().at['75%', i]
    IQR = Q3-Q1
    LTV = Q1 - 1.5 * IQR 
    UTV = Q3 + 1.5 * IQR 
    return LTV,UTV

# QUERY MAREAS ACODERE

driver = '{ODBC Driver 17 for SQL Server}'
server = 'srv-db-east-repl001.database.windows.net'
database_name = 'db_cfa_prd01'
user = 'userpowerbi'
password = '#p4ssw0rdp0w3rb1#'

connection_string = textwrap.dedent('''
    Driver={driver};
    Server={server};
    Database={database_name};
    Uid={user};
    Pwd={password};
    Encrypt=yes;
    TrustServerCertificate=no;
'''.format(driver=driver,server=server,database_name=database_name,user=user,password=password))

conn = pyodbc.connect(connection_string)

cursor = conn.cursor()

cursor.execute("""SELECT marea_id as MAREA,
       DATEADD(HOUR,-5,acodera_chata) AS ACODERA,
	     DATEADD(HOUR,-5,inicio_succion) AS INICIO_SUCCION,
	   DATEADD(HOUR,-5,termino_succion) AS TERMINO_SUCCION,
	   CASE WHEN DATEADD(HOUR,-5,acodera_chata) >= '{a}' AND DATEADD(HOUR,-5,acodera_chata) <= '{b}' THEN '2022 - I' 
	   WHEN DATEADD(HOUR,-5,acodera_chata) >= '{c}' AND DATEADD(HOUR,-5,acodera_chata) <= '{d}' THEN '2022 - II'
	   ELSE 'OTHER' END AS TEMPORADA
       FROM dbo.MareasAcodere
       WHERE DATEADD(HOUR,-5,acodera_chata) >= '{inicio}'""".format(a=inicio_tem_2022_1,b=fin_tem_2022_1,c=inicio_tem_2022_2,d=fin_tem_2022_2,inicio=inicio_tem_2022_1))

df_mareas = pd.DataFrame.from_records(cursor.fetchall(), columns=[col[0] for col in cursor.description])
df_mareas['MAREA'] = pd.to_numeric(df_mareas['MAREA'])
df_mareas = df_mareas[~df_mareas['MAREA'].duplicated()]
cursor.close()

driver = '{ODBC Driver 17 for SQL Server}'
server = 'srv-db-east-us-tasa-his-02.database.windows.net'
database_name = 'db_bi_production_prd'
user = 'userpowerbi'
password = '#p4ssw0rdp0w3rb1#'

connection_string = textwrap.dedent('''
    Driver={driver};
    Server={server};
    Database={database_name};
    Uid={user};
    Pwd={password};
    Encrypt=yes;
    TrustServerCertificate=no;
'''.format(driver=driver,server=server,database_name=database_name,user=user,password=password))

conn = pyodbc.connect(connection_string)

cursor = conn.cursor()

cursor.execute("""SELECT 
        MAREA, 
        PLANTA,
        CAST(F_ARRIB AS DATETIME) + CAST(H_ARRIB AS DATETIME) AS FECHA_ARRIBO,
        CASE WHEN D_INICI < H_ARRIB THEN CAST(DATEADD(DAY,1,F_ARRIB) AS DATETIME) + CAST(D_INICI AS DATETIME) ELSE CAST(F_ARRIB AS DATETIME) + CAST(D_INICI AS DATETIME) END AS INICIO_DESCARGA,
        CASE WHEN D_FINAL < D_INICI AND D_INICI < H_ARRIB THEN CAST(DATEADD(DAY,2,F_ARRIB) AS DATETIME) + CAST(D_FINAL AS DATETIME) 
         WHEN D_FINAL < D_INICI AND D_INICI > H_ARRIB THEN CAST(DATEADD(DAY,1,F_ARRIB) AS DATETIME) + CAST(D_FINAL AS DATETIME) 
         WHEN D_FINAL > D_INICI AND D_INICI < H_ARRIB THEN CAST(DATEADD(DAY,1,F_ARRIB) AS DATETIME) + CAST(D_FINAL AS DATETIME)
         ELSE CAST(F_ARRIB AS DATETIME) + CAST(D_FINAL AS DATETIME) END AS FIN_DESCARGA
        FROM sap.ZQM_DESC_MP_CHI
        WHERE FECHA >= '{inicio}' AND MAREA <> 0""".format(inicio=inicio_tem_2022_1))

df_chi = pd.DataFrame.from_records(cursor.fetchall(), columns=[col[0] for col in cursor.description])
df_chi['MAREA'] = pd.to_numeric(df_chi['MAREA'])
df_chi['PLANTA'] = df_chi['PLANTA'].str.split(' ').str[0].str.upper()
df_chi = df_chi[~df_chi['MAREA'].duplicated()]
cursor.close()

df_bitacora = pd.merge(df_mareas,df_chi, on=["MAREA"],how='inner')

# FORMS QUERY
df_forms_dgsante = pd.read_csv('https://docs.google.com/spreadsheets/d/e/2PACX-1vRsp2vmTU-CMxMpwWAui3SohhzlVSfQF-s5vGO3H8tkmFPjrKZk8ZZ47xflO-YQf4uMIkvVTiEX3aAe/pub?gid=0&single=true&output=csv')

df_forms_dgsante['1 - INICIO DE ANÁLISIS 1 -INICIO DE ANÁLISIS'] = pd.to_datetime(df_forms_dgsante['1 - INICIO DE ANÁLISIS 1 -INICIO DE ANÁLISIS'])
df_forms_dgsante['2 - FIN DE ANÁLISIS 2 - FIN DE ANÁLISIS'] = pd.to_datetime(df_forms_dgsante['2 - FIN DE ANÁLISIS 2 - FIN DE ANÁLISIS'])

df_forms_dgsante.rename(columns={'1 - INICIO DE ANÁLISIS 1 -INICIO DE ANÁLISIS':'INICIO_ANALISIS'},inplace = True)
df_forms_dgsante.rename(columns={'2 - FIN DE ANÁLISIS 2 - FIN DE ANÁLISIS':'FIN_ANALISIS'},inplace = True)

df_forms_dgsante['INICIO_ANALISIS'] = df_forms_dgsante['INICIO_ANALISIS'].apply(lambda x:  solve_datetime_sheets(x.to_pydatetime()))
df_forms_dgsante['FIN_ANALISIS'] = df_forms_dgsante['FIN_ANALISIS'].apply(lambda x:  solve_datetime_sheets(x.to_pydatetime()))

df_forms_dgsante = df_forms_dgsante[~df_forms_dgsante['MAREA'].duplicated()]

df_tasa = pd.merge(df_bitacora,df_forms_dgsante, on=["MAREA","PLANTA"],how='left')
df_tasa = df_tasa[~df_tasa['INICIO_SUCCION'].isnull()]
df_tasa = df_tasa[df_tasa['TEMPORADA']!='OTHER']

# DATA EXPORT EXCEL

#df_export = df_tasa.sort_values(by='ACODERA', ascending=False)
#df_export.to_excel(path_dgsante_project+"/consolidated_data.xlsx",index=False)

# ANALYTICS DATA 

df1 = df_tasa[['PLANTA','TEMPORADA','ACODERA','INICIO_SUCCION']]
df1['P'] = (df1['INICIO_SUCCION'] - df1['ACODERA'] ).dt.total_seconds() / 60.0
df1 = df1.groupby(['TEMPORADA'],as_index=False).agg({'P':'mean'})
df1.rename(columns = {'P':'X̅ P'}, inplace = True)
df1['X̅ P'] = df1['X̅ P'].round()
df1 = df1.sort_index()

barplot_df('TEMPORADA','X̅ P',df1,'a')


df1 = df_tasa[['PLANTA','TEMPORADA','ACODERA','INICIO_SUCCION']]
df1['P'] = (df1['INICIO_SUCCION'] - df1['ACODERA'] ).dt.total_seconds() / 60.0
df1 = df1.groupby(['TEMPORADA','PLANTA'],as_index=False).agg({'P':'mean'})
df1.rename(columns = {'P':'X̅ P'}, inplace = True)
df1['X̅ P'] = df1['X̅ P'].round()
df1 = df1.sort_index()

barplot_df_hue('PLANTA','X̅ P','TEMPORADA',df1,'b')

df_tasa['INISUCCION_INIDESCARGA_CORRECTO'] = np.where(df_tasa['INICIO_DESCARGA'] < df_tasa['INICIO_SUCCION'],0,1)
df_tasa['INISUCCION_INIDESCARGA_MIN'] = (df_tasa['INICIO_DESCARGA'] - df_tasa['INICIO_SUCCION'] ).dt.total_seconds() / 60.0

piechart_percentage(values_total('INISUCCION_INIDESCARGA_CORRECTO',df_tasa),'c')

LTV,UTV = select_values_outliers(df_tasa['INISUCCION_INIDESCARGA_MIN'].to_frame())
df_tasa['ATIPICO_INISUCCION_INIDESCARGA'] = np.where((df_tasa['INISUCCION_INIDESCARGA_MIN'] < LTV) | (df_tasa['INISUCCION_INIDESCARGA_MIN'] > UTV),1,0)

histplot_df_hue(df_tasa[df_tasa['ATIPICO_INISUCCION_INIDESCARGA']==0],'INISUCCION_INIDESCARGA_MIN','TEMPORADA','d')

histplot_df_hue(df_tasa[ (df_tasa['ATIPICO_INISUCCION_INIDESCARGA']==0) & (df_tasa['TEMPORADA'] == '2022 - I')].sort_values(by='PLANTA'),'INISUCCION_INIDESCARGA_MIN','PLANTA','e')

histplot_df_hue(df_tasa[ (df_tasa['ATIPICO_INISUCCION_INIDESCARGA']==0) & (df_tasa['TEMPORADA'] == '2022 - II')].sort_values(by='PLANTA'),'INISUCCION_INIDESCARGA_MIN','PLANTA','f')

df_tasa2 = df_tasa[~df_tasa['INICIO_ANALISIS'].isnull()]
df_tasa2 = df_tasa2[df_tasa2['ATIPICO_INISUCCION_INIDESCARGA']==0]

df_tasa2['FINANALISIS_INISUCCION_CORRECTO'] = np.where(df_tasa2['FIN_ANALISIS'] > df_tasa2['INICIO_SUCCION'],0,1) 
df_tasa2['FINANALISIS_INISUCCION_MIN'] = (df_tasa2['INICIO_SUCCION'] - df_tasa2['FIN_ANALISIS'] ).dt.total_seconds() / 60.0

piechart_percentage(values_total('FINANALISIS_INISUCCION_CORRECTO',df_tasa2),'g')

LTV,UTV = select_values_outliers(df_tasa2['FINANALISIS_INISUCCION_MIN'].to_frame())
df_tasa2['ATIPICO_FINANALISIS_INISUCCION'] = np.where((df_tasa2['FINANALISIS_INISUCCION_MIN'] < LTV) | (df_tasa2['FINANALISIS_INISUCCION_MIN'] > UTV),1,0)

df1 = df_tasa2[df_tasa2['ATIPICO_FINANALISIS_INISUCCION'] == 0][['TEMPORADA','PLANTA','FIN_ANALISIS','INICIO_SUCCION']]
df1['K'] = (df1['INICIO_SUCCION'] - df1['FIN_ANALISIS'] ).dt.total_seconds() / 60.0
df1 = df1.groupby(['TEMPORADA'],as_index=False).agg({'K':'mean'})
df1.rename(columns = {'K':'X̅ K'}, inplace = True)
df1['X̅ K'] = df1['X̅ K'].round()
df1 = df1.sort_index()

barplot_df('TEMPORADA','X̅ K',df1,'h')

df1 = df_tasa2[df_tasa2['ATIPICO_FINANALISIS_INISUCCION'] == 0][['PLANTA','TEMPORADA','FIN_ANALISIS','INICIO_SUCCION']]
df1['K'] = (df1['INICIO_SUCCION'] - df1['FIN_ANALISIS'] ).dt.total_seconds() / 60.0
df1 = df1.groupby(['TEMPORADA','PLANTA'],as_index=False).agg({'K':'mean'})
df1.rename(columns = {'K':'X̅ K'}, inplace = True)
df1['X̅ K'] = df1['X̅ K'].round()
df1 = df1.sort_index()

barplot_df_hue('PLANTA','X̅ K','TEMPORADA',df1,'i')

df_tasa2 = df_tasa[~df_tasa['INICIO_ANALISIS'].isnull()]
df_tasa2 = df_tasa2[df_tasa2['ATIPICO_INISUCCION_INIDESCARGA']==0]

df_tasa2['INIANALISIS_INISUCCION_CORRECTO'] = np.where(df_tasa2['INICIO_ANALISIS'] > df_tasa2['INICIO_SUCCION'],0,1) 
df_tasa2['INIANALISIS_INISUCCION_MIN'] = (df_tasa2['INICIO_SUCCION'] - df_tasa2['INICIO_ANALISIS'] ).dt.total_seconds() / 60.0

piechart_percentage(values_total('INIANALISIS_INISUCCION_CORRECTO',df_tasa2),'j')

LTV,UTV = select_values_outliers(df_tasa2['INIANALISIS_INISUCCION_MIN'].to_frame())
df_tasa2['ATIPICO_INIANALISIS_INISUCCION'] = np.where((df_tasa2['INIANALISIS_INISUCCION_MIN'] < LTV) | (df_tasa2['INIANALISIS_INISUCCION_MIN'] > UTV),1,0)

df1 = df_tasa2[df_tasa2['ATIPICO_INIANALISIS_INISUCCION'] == 0][['TEMPORADA','PLANTA','INICIO_ANALISIS','INICIO_SUCCION']]
df1['G'] = (df1['INICIO_SUCCION'] - df1['INICIO_ANALISIS'] ).dt.total_seconds() / 60.0
df1 = df1.groupby(['TEMPORADA'],as_index=False).agg({'G':'mean'})
df1.rename(columns = {'G':'X̅ G'}, inplace = True)
df1['X̅ G'] = df1['X̅ G'].round()
df1 = df1.sort_index()

barplot_df('TEMPORADA','X̅ G',df1,'k')

df1 = df_tasa2[df_tasa2['ATIPICO_INIANALISIS_INISUCCION'] == 0][['PLANTA','TEMPORADA','INICIO_ANALISIS','INICIO_SUCCION']]
df1['G'] = (df1['INICIO_SUCCION'] - df1['INICIO_ANALISIS'] ).dt.total_seconds() / 60.0
df1 = df1.groupby(['TEMPORADA','PLANTA'],as_index=False).agg({'G':'mean'})
df1.rename(columns = {'G':'X̅ G'}, inplace = True)
df1['X̅ G'] = df1['X̅ G'].round()
df1 = df1.sort_index()

barplot_df_hue('PLANTA','X̅ G','TEMPORADA',df1,'l')

# UPDATE PPT

pptx = Presentation(template_pptx)

slide_acodera_inisuccion1 = pptx.slides[2]
slide_acodera_inisuccion2 = pptx.slides[3]

slide_inisuccion_inidescarga1 = pptx.slides[5]
slide_inisuccion_inidescarga2 = pptx.slides[6]
slide_inisuccion_inidescarga3 = pptx.slides[7]
slide_inisuccion_inidescarga4 = pptx.slides[8]

slide_finanalisis_inisuccion1 = pptx.slides[10]
slide_finanalisis_inisuccion2 = pptx.slides[11]
slide_finanalisis_inisuccion3 = pptx.slides[12]

slide_inianalisis_inisuccion1 = pptx.slides[14]
slide_inianalisis_inisuccion2 = pptx.slides[15]
slide_inianalisis_inisuccion3 = pptx.slides[16]


left = top = Inches(1)
slide_acodera_inisuccion1.shapes.add_picture(path_resources+'/'+'a.png',left,top)
slide_acodera_inisuccion2.shapes.add_picture(path_resources+'/'+'b.png',left,top)
slide_inisuccion_inidescarga1.shapes.add_picture(path_resources+'/'+'c.png',left,top)
slide_inisuccion_inidescarga2.shapes.add_picture(path_resources+'/'+'d.png',left,top)
slide_inisuccion_inidescarga3.shapes.add_picture(path_resources+'/'+'e.png',left,top)
slide_inisuccion_inidescarga4.shapes.add_picture(path_resources+'/'+'f.png',left,top)
slide_finanalisis_inisuccion1.shapes.add_picture(path_resources+'/'+'g.png',left,top)
slide_finanalisis_inisuccion2.shapes.add_picture(path_resources+'/'+'h.png',left,top)
slide_finanalisis_inisuccion3.shapes.add_picture(path_resources+'/'+'i.png',left,top)
slide_inianalisis_inisuccion1.shapes.add_picture(path_resources+'/'+'j.png',left,top)
slide_inianalisis_inisuccion2.shapes.add_picture(path_resources+'/'+'k.png',left,top)
slide_inianalisis_inisuccion3.shapes.add_picture(path_resources+'/'+'l.png',left,top)


pptx.save(path_dgsante_project+'/resultados_dgsante.pptx')
